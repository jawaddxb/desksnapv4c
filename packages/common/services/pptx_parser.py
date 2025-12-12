"""
PPTX Parser Service
Extracts slides from PowerPoint files using python-pptx
KISS: Simple extraction logic with clear output structure
"""
import base64
import io
import logging
from typing import Any

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str) -> list[dict[str, Any]]:
    """
    Extract slides from PPTX file.

    Returns list of raw slide data:
    - texts: [{value, font_size, font_family, is_bold, position}]
    - images: [{data_url, width, height}]
    - notes: speaker notes text
    - background_color: hex color
    """
    try:
        prs = Presentation(file_path)
        slides_data = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                "position": slide_idx,
                "texts": [],
                "images": [],
                "notes": get_slide_notes(slide),
                "background_color": None,
            }

            # Extract content from shapes
            for shape in slide.shapes:
                # Text extraction
                text_data = extract_text_from_shape(shape)
                if text_data:
                    slide_data["texts"].append(text_data)

                # Image extraction
                image_data = extract_image_from_shape(shape)
                if image_data:
                    slide_data["images"].append(image_data)

            slides_data.append(slide_data)

        return slides_data

    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        raise ValueError(f"Failed to parse PPTX file: {str(e)}")


def extract_text_from_shape(shape: BaseShape) -> dict[str, Any] | None:
    """Extract text with formatting from a shape."""
    if not shape.has_text_frame:
        return None

    # Get all text from the shape
    full_text = ""
    font_sizes: list[float] = []
    is_bold = False
    font_family = "Arial"

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            full_text += run.text
            if run.font.size:
                font_sizes.append(run.font.size.pt)
            if run.font.bold:
                is_bold = True
            if run.font.name:
                font_family = run.font.name

    full_text = full_text.strip()
    if not full_text:
        return None

    # Calculate average font size
    avg_font_size = sum(font_sizes) / len(font_sizes) if font_sizes else 12

    # Get position
    left = shape.left.pt if shape.left else 0
    top = shape.top.pt if shape.top else 0
    width = shape.width.pt if shape.width else 0
    height = shape.height.pt if shape.height else 0

    # Determine if this is likely a title
    is_title = (
        avg_font_size > 24 and
        top < 150 and  # Near top of slide
        len(full_text) < 200  # Not too long
    )

    return {
        "value": full_text,
        "font_size": avg_font_size,
        "font_family": font_family,
        "is_bold": is_bold,
        "is_title": is_title,
        "position": {
            "left": left,
            "top": top,
            "width": width,
            "height": height,
        },
    }


def extract_image_from_shape(shape: BaseShape) -> dict[str, Any] | None:
    """Extract image as base64 data URL."""
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return None

        image = shape.image
        image_bytes = image.blob
        content_type = image.content_type

        # Convert to base64 data URL
        base64_data = base64.b64encode(image_bytes).decode("utf-8")
        data_url = f"data:{content_type};base64,{base64_data}"

        # Get dimensions
        width = shape.width.pt if shape.width else 0
        height = shape.height.pt if shape.height else 0

        return {
            "data_url": data_url,
            "width": width,
            "height": height,
            "content_type": content_type,
        }
    except Exception as e:
        logger.warning(f"Could not extract image: {e}")
        return None


def get_slide_notes(slide: Any) -> str:
    """Extract speaker notes from slide."""
    try:
        if not slide.has_notes_slide:
            return ""

        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        return notes_text_frame.text.strip() if notes_text_frame else ""
    except Exception as e:
        logger.warning(f"Could not extract notes: {e}")
        return ""


def get_slide_count(file_path: str) -> int:
    """Get the number of slides in a PPTX file."""
    try:
        prs = Presentation(file_path)
        return len(prs.slides)
    except Exception as e:
        logger.error(f"Error counting slides: {e}")
        return 0
